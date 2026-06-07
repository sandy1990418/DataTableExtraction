from __future__ import annotations

import io
import json
from types import SimpleNamespace

from fastapi.testclient import TestClient
from pptx import Presentation
import pytest

from app.main import create_app
from app.config import Settings
from app.models import EvidenceItem
from app.services.document_parser import parse_markdown
from app.services.canonical_extractor import _planning_evidence_to_prompt, plan_tables
from app.services.evidence_selector import select_table_evidence_blocks
from app.services.llm_service import (
    TableSpec,
    _merge_source_and_generated_tables,
    _source_tables_from_message,
    _table_catalog_summary,
    _tables_from_tool_args,
    chat,
)
from app.services.pipeline import validate_table_quality
from app.services.grounding import check_numeric_grounding, ground_table_reviews
from app.services.report_agent_team import (
    report_evidence_breakdown,
    run_table_report_agent_team,
    select_report_evidence_items,
)
from app.services.text_match import keywords, overlap_score
from app.services.table_agent_team import run_table_agent_team
from app.services.table_qa import review_tables
from app.services.table_extraction import extract_source_tables
from app.tools import table_pptx


def test_build_pptx_rejects_empty_headers() -> None:
    with pytest.raises(ValueError, match="headers must not be empty"):
        table_pptx.build_pptx("Title", [], [["value"]])


def test_table_spec_normalizes_rows_and_layout() -> None:
    table = TableSpec(
        title="  Quarterly Review  ",
        headers=["Name", "Status"],
        rows=[["Alice"]],
        text="   ",
        layout="text_left",
    )

    assert table.title == "Quarterly Review"
    assert table.rows == [["Alice", ""]]
    assert table.layout == "table_only"
    assert table.table_ratio == 0.5


def test_download_uses_stored_filename(tmp_path) -> None:
    original_store_dir = table_pptx.STORE_DIR
    table_pptx.STORE_DIR = tmp_path
    try:
        token = table_pptx.store_pptx(b"pptx-bytes", "中文 / 報表?.pptx", 60)
        client = TestClient(create_app())

        response = client.get(f"/download/{token}")

        assert response.status_code == 200
        assert response.content == b"pptx-bytes"
        assert response.headers["content-disposition"] == (
            'attachment; filename="table.pptx"; '
            "filename*=UTF-8''%E4%B8%AD%E6%96%87%20%E5%A0%B1%E8%A1%A8%3F.pptx"
        )
    finally:
        table_pptx.STORE_DIR = original_store_dir


def test_safe_filename_preserves_unicode() -> None:
    assert table_pptx.safe_pptx_filename("中文 / 報表?.pptx") == "中文 報表?.pptx"
    assert table_pptx.ascii_fallback_filename("中文 / 報表?.pptx") == "table.pptx"


def test_multi_table_args_preserve_source_table_before_derived_table() -> None:
    tables = _tables_from_tool_args(
        {
            "tables": [
                {
                    "title": "Original Metrics",
                    "kind": "source_table",
                    "headers": ["Metric", "Value"],
                    "rows": [["Revenue", "$10M"]],
                },
                {
                    "title": "Key Takeaways",
                    "kind": "extracted_summary",
                    "headers": ["Point", "Evidence"],
                    "rows": [["Revenue grew", "Revenue = $10M"]],
                    "text": "Derived from the surrounding prose.",
                    "layout": "table_bottom",
                    "table_ratio": 0.33,
                },
            ]
        }
    )

    assert [table.kind for table in tables] == ["source_table", "extracted_summary"]
    assert tables[0].rows == [["Revenue", "$10M"]]
    assert tables[1].layout == "table_bottom"
    assert tables[1].table_ratio == 0.33


def test_build_tables_pptx_creates_slide_per_table() -> None:
    pptx_bytes = table_pptx.build_tables_pptx(
        [
            {
                "title": "Original Table",
                "kind": "source_table",
                "headers": ["A", "B"],
                "rows": [["1", "2"]],
            },
            {
                "title": "Summary Table",
                "kind": "extracted_summary",
                "headers": ["Theme", "Detail"],
                "rows": [["Risk", "Needs follow-up"]],
            },
        ]
    )

    presentation = Presentation(io.BytesIO(pptx_bytes))
    assert len(presentation.slides) == 2


def test_table_bottom_ratio_places_table_in_lower_third() -> None:
    pptx_bytes = table_pptx.build_pptx(
        "Lower Third",
        ["A", "B"],
        [["1", "2"], ["3", "4"], ["5", "6"]],
        text="Context above the table.",
        layout="table_bottom",
        table_ratio=0.33,
    )

    presentation = Presentation(io.BytesIO(pptx_bytes))
    slide = presentation.slides[0]
    table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))

    assert table_shape.top > table_pptx.CONTENT_TOP + table_pptx.CONTENT_H * 0.5
    assert table_shape.height <= table_pptx.CONTENT_H * 0.33


def test_build_plan_pptx_preserves_non_table_slides() -> None:
    pptx_bytes = table_pptx.build_plan_pptx(
        {
            "presentation_title": "Demo",
            "slides": [
                {
                    "slide_number": 1,
                    "slide_type": "title",
                    "title": "Opening",
                    "content": "Context",
                    "table_ref": None,
                },
                {
                    "slide_number": 2,
                    "slide_type": "table",
                    "title": "Metrics",
                    "content": "Reported metrics",
                    "table_ref": "metrics",
                },
                {
                    "slide_number": 3,
                    "slide_type": "conclusion",
                    "title": "Conclusion",
                    "content": "Done",
                    "table_ref": None,
                },
            ],
        },
        [
            {
                "table_id": "metrics",
                "title": "Metrics Table",
                "headers": ["System", "F1"],
                "rows": [["MemoryOS", "49.11%"]],
            }
        ],
    )

    presentation = Presentation(io.BytesIO(pptx_bytes))
    assert len(presentation.slides) == 3
    assert [
        sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
        for slide in presentation.slides
    ] == [0, 1, 0]


def test_pdf_like_markdown_is_split_into_table_sections() -> None:
    doc = parse_markdown(
        """Memory OS of AI Agent
1
Introduction
Background text.

Table 2: LoCoMo dataset comparison with per-category scores and average ranks.
Model
Method
F1 ↑
BLEU-1 ↑
GPT-4o-mini
MemoryOS
49.11
46.18
2
Related Work
More prose.
""",
        "paper.md",
    )

    table_sections = [section for section in doc.sections if section.heading.startswith("Table 2")]
    assert len(table_sections) == 1
    assert "MemoryOS" in table_sections[0].content
    assert "49.11" in table_sections[0].content


def test_evidence_selector_prioritizes_matching_table_sections() -> None:
    items = [
        EvidenceItem(
            kind="text_fact",
            source_ref="paper:1",
            title="Introduction",
            content="General memory overview.",
        ),
        EvidenceItem(
            kind="text_fact",
            source_ref="paper:table:2",
            title="Table 2: LoCoMo dataset comparison",
            content="MemoryOS reports F1 49.11% and BLEU-1 46.18% on GPT-4o-mini.",
        ),
    ]
    blocks = select_table_evidence_blocks(
        [
            {
                "name": "locomo_metrics",
                "title": "LoCoMo Metrics",
                "description": "Compare F1 and BLEU-1 scores",
                "columns": [{"name": "F1"}, {"name": "BLEU-1"}],
            }
        ],
        items,
    )

    assert "Table 2: LoCoMo dataset comparison" in blocks["locomo_metrics"]
    assert "49.11%" in blocks["locomo_metrics"]


def test_table_planning_prompt_keeps_metric_rows_visible() -> None:
    prompt = _planning_evidence_to_prompt(
        [
            EvidenceItem(
                kind="text_fact",
                source_ref="paper:lines:1-2",
                title="Introduction",
                content="General background on memory systems.",
            ),
            EvidenceItem(
                kind="text_fact",
                source_ref="paper:lines:741-771",
                title="Table 2: LoCoMo dataset comparison",
                content=(
                    "Model Method Single Hop Multi Hop Temporal Open Domain Avg. Rank "
                    "F1 BLEU-1 GPT-4o-mini MemoryOS 49.11 46.18 2.0 2.0"
                ),
            ),
        ]
    )

    assert "Detailed evidence for table planning" in prompt
    assert "Table 2: LoCoMo dataset comparison" in prompt
    assert "MemoryOS 49.11 46.18" in prompt


async def test_plan_tables_sends_detailed_table_evidence(monkeypatch) -> None:
    captured_prompts = []

    class FakeCompletions:
        async def create(self, **kwargs):
            captured_prompts.append(kwargs["messages"][1]["content"])
            return SimpleNamespace(
                choices=[
                    SimpleNamespace(
                        message=SimpleNamespace(
                            content=json.dumps(
                                {
                                    "tables": [
                                        {
                                            "name": "locomo_metrics",
                                            "title": "LoCoMo Metrics",
                                            "description": "Compare exact LoCoMo metrics.",
                                            "row_entity": "one method",
                                            "evidence_anchors": [
                                                "Table 2: LoCoMo dataset comparison"
                                            ],
                                            "columns": [
                                                {
                                                    "name": "Method",
                                                    "description": "method name",
                                                    "example": "MemoryOS",
                                                },
                                                {
                                                    "name": "F1",
                                                    "description": "reported F1 score",
                                                    "example": "49.11",
                                                },
                                            ],
                                        }
                                    ]
                                }
                            )
                        )
                    )
                ]
            )

    class FakeAsyncOpenAI:
        def __init__(self, **_kwargs):
            self.chat = SimpleNamespace(completions=FakeCompletions())

    monkeypatch.setattr("app.services.canonical_extractor.AsyncOpenAI", FakeAsyncOpenAI)

    specs, _evidence_block = await plan_tables(
        [
            EvidenceItem(
                kind="text_fact",
                source_ref="paper:lines:741-771",
                title="Table 2: LoCoMo dataset comparison",
                content="MemoryOS reports F1 49.11 and BLEU-1 46.18 on GPT-4o-mini.",
            )
        ],
        Settings(OPENAI_API_KEY="test"),
    )

    assert specs[0]["name"] == "locomo_metrics"
    assert "MemoryOS reports F1 49.11" in captured_prompts[0]


def test_validate_table_quality_flags_sparse_and_generic_tables() -> None:
    warnings = validate_table_quality(
        [
            {"table_id": "sparse", "rows": [["-", "N/A"], ["", "-"]]},
            {
                "table_id": "generic",
                "rows": [["Yes", "High"], ["Moderate", "Low"], ["Yes", "High"]],
            },
        ]
    )

    assert any("sparse" in warning for warning in warnings)
    assert any("generic" in warning for warning in warnings)


def test_report_evidence_agent_prioritizes_hint_relevant_tables() -> None:
    items = [
        EvidenceItem(
            kind="text_fact",
            source_ref="paper:intro",
            title="Introduction",
            content="General discussion about agents.",
        ),
        EvidenceItem(
            kind="text_fact",
            source_ref="paper:table:2",
            title="Table 2: LoCoMo dataset comparison",
            content="MemoryOS reports F1 49.11 and BLEU-1 46.18 on GPT-4o-mini.",
        ),
        EvidenceItem(
            kind="text_fact",
            source_ref="paper:method",
            title="Memory Architecture",
            content="MemoryOS uses short-term, mid-term, and long-term memory storage.",
        ),
    ]

    selected = select_report_evidence_items(items, "compare LoCoMo F1 BLEU memory architecture")

    assert selected[0].title == "Table 2: LoCoMo dataset comparison"
    assert {item.title for item in selected[:2]} == {
        "Table 2: LoCoMo dataset comparison",
        "Memory Architecture",
    }


def test_extract_source_tables_builds_catalog_from_markdown() -> None:
    tables = extract_source_tables(
        """Quarterly Revenue

| Product | Revenue | Growth |
| --- | ---: | ---: |
| A | 120 | 20% |
| B | 80 | -5% |
"""
    )

    assert len(tables) == 1
    assert tables[0]["table_id"] == "tbl_001"
    assert tables[0]["kind"] == "source_table"
    assert tables[0]["title"] == "Quarterly Revenue"
    assert tables[0]["headers"] == ["Product", "Revenue", "Growth"]
    assert tables[0]["rows"] == [["A", "120", "20%"], ["B", "80", "-5%"]]
    assert tables[0]["summary"] == "2 rows x 3 columns (Product, Revenue, Growth)"
    assert tables[0]["source_ref"] == "lines:3-6"


def test_extract_source_tables_builds_catalog_from_html() -> None:
    tables = extract_source_tables(
        """
<table>
  <caption>Pipeline Status</caption>
  <tr><th>Stage</th><th>Status</th></tr>
  <tr><td>Outline</td><td>Done</td></tr>
</table>
"""
    )

    assert len(tables) == 1
    assert tables[0]["table_id"] == "tbl_001"
    assert tables[0]["title"] == "Pipeline Status"
    assert tables[0]["headers"] == ["Stage", "Status"]
    assert tables[0]["rows"] == [["Outline", "Done"]]
    assert tables[0]["source_ref"] == "html:table"


def test_source_table_catalog_is_preserved_before_generated_tables() -> None:
    source_tables = _source_tables_from_message(
        """Original Data

| Name | Score |
| --- | --- |
| Alice | 95 |
"""
    )
    generated_tables = _tables_from_tool_args(
        {
            "tables": [
                {
                    "title": "Duplicate Source",
                    "kind": "source_table",
                    "headers": ["Name", "Score"],
                    "rows": [["Alice", "95"]],
                },
                {
                    "title": "Derived Summary",
                    "kind": "extracted_summary",
                    "headers": ["Finding", "Evidence"],
                    "rows": [["Alice leads", "Score 95"]],
                },
            ]
        }
    )

    merged = _merge_source_and_generated_tables(source_tables, generated_tables)

    assert [table.kind for table in merged] == ["source_table", "extracted_summary"]
    assert merged[0].table_id == "tbl_001"
    assert merged[0].title == "Original Data"
    assert "tbl_001: Original Data" in _table_catalog_summary(source_tables)


async def test_chat_returns_error_for_invalid_tool_table_args(monkeypatch) -> None:
    invalid_args = {
        "tables": [
            {
                "title": "Broken Table",
                "kind": "extracted_summary",
                "headers": ["Name", ""],
                "rows": [["Alice", "95"]],
            }
        ]
    }

    class FakeCompletions:
        async def create(self, **_kwargs):
            return SimpleNamespace(
                choices=[
                    SimpleNamespace(
                        message=SimpleNamespace(
                            content=None,
                            tool_calls=[
                                SimpleNamespace(
                                    function=SimpleNamespace(arguments=json.dumps(invalid_args))
                                )
                            ],
                        )
                    )
                ]
            )

    class FakeAsyncOpenAI:
        def __init__(self, **_kwargs):
            self.chat = SimpleNamespace(completions=FakeCompletions())

    monkeypatch.setattr("app.services.llm_service.AsyncOpenAI", FakeAsyncOpenAI)

    response = await chat("make a table", Settings(OPENAI_API_KEY="test"))

    assert response["type"] == "error"
    assert response["message"].startswith("Tool args validation error:")


async def test_review_tables_uses_table_specific_evidence(monkeypatch) -> None:
    captured_user_messages = []

    class FakeCompletions:
        async def create(self, **kwargs):
            captured_user_messages.append(kwargs["messages"][1]["content"])
            return SimpleNamespace(
                choices=[
                    SimpleNamespace(
                        message=SimpleNamespace(
                            content=json.dumps(
                                {
                                    "table_id": "metrics",
                                    "status": "pass",
                                    "warnings": [],
                                    "unsupported_cells": [],
                                }
                            )
                        )
                    )
                ]
            )

    class FakeAsyncOpenAI:
        def __init__(self, **_kwargs):
            self.chat = SimpleNamespace(completions=FakeCompletions())

    monkeypatch.setattr("app.services.table_qa.AsyncOpenAI", FakeAsyncOpenAI)

    reviews = await review_tables(
        [{"table_id": "metrics", "title": "Metrics", "headers": ["Metric"], "rows": [["49.11%"]]}],
        {"metrics": "focused evidence with 49.11%", "": "fallback evidence"},
        Settings(OPENAI_API_KEY="test"),
    )

    assert reviews[0]["status"] == "pass"
    assert "focused evidence with 49.11%" in captured_user_messages[0]
    assert "fallback evidence" not in captured_user_messages[0]


async def test_table_agent_team_runs_revision_loop(monkeypatch) -> None:
    async def fake_populate_tables(_specs, _evidence_block, _settings):
        return [
            {
                "name": "metrics",
                "title": "Metrics",
                "description": "",
                "headers": ["System", "F1"],
                "rows": [["MemoryOS", "High"]],
            }
        ]

    review_calls = 0

    async def fake_review_tables(tables, _evidence_block, _settings):
        nonlocal review_calls
        review_calls += 1
        if review_calls == 1:
            return [
                {
                    "table_id": tables[0]["table_id"],
                    "status": "needs_revision",
                    "warnings": ["Generic qualitative value is unsupported."],
                    "unsupported_cells": [
                        {
                            "row": 1,
                            "column": "F1",
                            "value": "High",
                            "reason": "Evidence reports 49.11%, not High.",
                        }
                    ],
                }
            ]
        return [
            {
                "table_id": tables[0]["table_id"],
                "status": "pass",
                "warnings": [],
                "unsupported_cells": [],
            }
        ]

    async def fake_revise_tables(tables, _reviews, _evidence_block, _settings):
        return [{**tables[0], "rows": [["MemoryOS", "49.11%"]]}]

    monkeypatch.setattr("app.services.table_agent_team.populate_tables", fake_populate_tables)
    monkeypatch.setattr("app.services.table_agent_team.review_tables", fake_review_tables)
    monkeypatch.setattr("app.services.table_agent_team.revise_tables", fake_revise_tables)

    result = await run_table_agent_team(
        {
            "slides": [
                {
                    "slide_number": 1,
                    "slide_type": "table",
                    "title": "Metrics",
                    "table_ref": "metrics",
                }
            ]
        },
        [
            {
                "name": "metrics",
                "title": "Metrics",
                "description": "Compare reported F1 scores",
                "columns": [{"name": "System"}, {"name": "F1"}],
            }
        ],
        [
            EvidenceItem(
                kind="text_fact",
                source_ref="paper:table:2",
                title="Table 2: LoCoMo dataset comparison",
                content="MemoryOS reports F1 49.11% on GPT-4o-mini.",
            )
        ],
        "fallback evidence",
        Settings(OPENAI_API_KEY="test"),
    )

    assert result.tables[0]["rows"] == [["MemoryOS", "49.11%"]]
    assert result.qa_reviews[0]["status"] == "pass"
    assert [event["agent"] for event in result.agent_trace] == [
        "SupervisorAgent",
        "EvidenceSelectorAgent",
        "TablePopulatorAgent",
        "TableReviewerAgent",
        "TableRevisionAgent",
        "TableReviewerAgent",
    ]


async def test_table_report_agent_team_reviews_specs_before_population(monkeypatch) -> None:
    call_order = []

    async def fake_plan_tables(_items, _settings, hint=""):
        call_order.append("plan_specs")
        return [
            {
                "name": "generic_metrics",
                "title": "Generic Metrics",
                "description": "Bad generic metric table",
                "row_entity": "one method",
                "columns": [
                    {"name": "Method", "description": "method name", "example": "MemoryOS"},
                    {"name": "Accuracy", "description": "generic score", "example": "49.11"},
                ],
            }
        ], "planning evidence with F1 49.11 and BLEU-1 46.18"

    spec_review_calls = 0

    async def fake_review_table_specs(specs, _hint, _evidence_block, _settings):
        nonlocal spec_review_calls
        call_order.append("review_specs")
        spec_review_calls += 1
        if spec_review_calls == 1:
            return [
                {
                    "name": specs[0]["name"],
                    "status": "needs_revision",
                    "warnings": ["Use exact metric names from evidence."],
                    "unsupported_columns": [
                        {"column": "Accuracy", "reason": "Evidence reports F1 and BLEU-1."}
                    ],
                }
            ]
        return [
            {
                "name": specs[0]["name"],
                "status": "pass",
                "warnings": [],
                "unsupported_columns": [],
            }
        ]

    async def fake_revise_table_specs(_specs, _reviews, _hint, _evidence_block, _settings):
        call_order.append("revise_specs")
        return [
            {
                "name": "locomo_metrics",
                "title": "LoCoMo Metrics",
                "description": "Compare exact LoCoMo F1 and BLEU-1 metrics.",
                "row_entity": "one method",
                "evidence_anchors": ["Table 2: LoCoMo dataset comparison"],
                "columns": [
                    {"name": "Method", "description": "method name", "example": "MemoryOS"},
                    {"name": "F1", "description": "reported F1", "example": "49.11"},
                    {"name": "BLEU-1", "description": "reported BLEU-1", "example": "46.18"},
                ],
            }
        ]

    async def fake_populate_tables(specs, _evidence_block, _settings):
        call_order.append("populate_rows")
        assert specs[0]["name"] == "locomo_metrics"
        assert [column["name"] for column in specs[0]["columns"]] == ["Method", "F1", "BLEU-1"]
        return [
            {
                "name": "locomo_metrics",
                "title": "LoCoMo Metrics",
                "description": "Compare exact LoCoMo F1 and BLEU-1 metrics.",
                "headers": ["Method", "F1", "BLEU-1"],
                "rows": [["MemoryOS", "49.11", "46.18"]],
            }
        ]

    async def fake_review_tables(tables, _evidence_block, _settings):
        call_order.append("review_rows")
        return [
            {
                "table_id": tables[0]["table_id"],
                "status": "pass",
                "warnings": [],
                "unsupported_cells": [],
            }
        ]

    monkeypatch.setattr("app.services.report_agent_team.plan_tables", fake_plan_tables)
    monkeypatch.setattr(
        "app.services.report_agent_team.review_table_specs", fake_review_table_specs
    )
    monkeypatch.setattr(
        "app.services.report_agent_team.revise_table_specs", fake_revise_table_specs
    )
    monkeypatch.setattr("app.services.report_agent_team.populate_tables", fake_populate_tables)
    monkeypatch.setattr("app.services.report_agent_team.review_tables", fake_review_tables)

    result = await run_table_report_agent_team(
        [
            EvidenceItem(
                kind="text_fact",
                source_ref="paper:table:2",
                title="Table 2: LoCoMo dataset comparison",
                content="MemoryOS reports F1 49.11 and BLEU-1 46.18.",
            )
        ],
        Settings(OPENAI_API_KEY="test"),
        hint="compare LoCoMo F1 BLEU metrics",
        n_slides=5,
    )

    assert call_order == [
        "plan_specs",
        "review_specs",
        "revise_specs",
        "review_specs",
        "populate_rows",
        "review_rows",
    ]
    assert result.specs[0]["name"] == "locomo_metrics"
    assert result.tables[0]["rows"] == [["MemoryOS", "49.11", "46.18"]]
    assert result.ppt_plan["slides"][2]["table_ref"] == "locomo_metrics"
    assert [event["agent"] for event in result.agent_trace] == [
        "ReportSupervisorAgent",
        "ReportEvidenceAgent",
        "TableIntentAgent",
        "TableSpecReviewAgent",
        "TableSpecRevisionAgent",
        "TableSpecReviewAgent",
        "TableEvidenceAgent",
        "TablePopulatorAgent",
        "TableGroundingAgent",
        "ReportComposerAgent",
    ]


def test_keywords_tokenizes_chinese_hint() -> None:
    # ASCII-only tokenization dropped CJK hints entirely; the shared tokenizer
    # must yield CJK unigrams + bigrams so overlap with CJK evidence is non-zero.
    query = keywords("比較 ASR 錯誤率")
    assert "asr" in query
    assert "錯誤" in query  # bigram
    assert "錯" in query  # unigram

    evidence = keywords("各模型的語音辨識錯誤率比較 (ASR error rate)")
    assert overlap_score(query, evidence) > 0


def test_report_evidence_breakdown_exposes_scores_and_drops() -> None:
    items = [
        EvidenceItem(
            kind="markdown_table",
            source_ref="paper:table:1",
            title="Table 1: Latency",
            content="3 rows x 2 columns",
            headers=["Model", "Latency"],
            rows=[["A", "12ms"]],
        ),
        EvidenceItem(
            kind="text_fact",
            source_ref="paper:intro",
            title="Introduction",
            content="Unrelated prose.",
        ),
    ]

    breakdown = report_evidence_breakdown(items, "latency comparison", max_items=1)

    assert breakdown["query_terms"]
    assert breakdown["selected"][0]["source_ref"] == "paper:table:1"
    assert "score" in breakdown["selected"][0]
    # The table that didn't make the cut is surfaced as dropped high-value evidence
    # so a retrieval miss is visible in the trace instead of guessed at.
    assert isinstance(breakdown["dropped_high_value_evidence"], list)


def test_numeric_grounding_flags_fabricated_numbers() -> None:
    table = {
        "table_id": "metrics",
        "headers": ["Model", "F1"],
        "rows": [["A", "49.11"], ["B", "99.99"]],
    }
    evidence = "Model A reports F1 49.11 on the benchmark."

    unsupported = check_numeric_grounding(table, evidence)

    assert len(unsupported) == 1
    assert unsupported[0]["number"] == "99.99"
    assert unsupported[0]["row"] == 1


def test_ground_table_reviews_forces_revision_on_ungrounded_cell() -> None:
    tables = [
        {
            "table_id": "metrics",
            "headers": ["Model", "F1"],
            "rows": [["A", "12.34"]],
        }
    ]
    reviews = [{"table_id": "metrics", "status": "pass", "warnings": [], "unsupported_cells": []}]

    grounded = ground_table_reviews(tables, {"metrics": "no numbers here at all"}, reviews)

    assert grounded[0]["status"] == "needs_revision"
    assert grounded[0]["unsupported_cells"]
    # Original review object is not mutated.
    assert reviews[0]["status"] == "pass"


def test_parser_picks_up_html_tables_missed_by_pipe_scanner() -> None:
    doc = parse_markdown(
        """# Results

<table>
<caption>HTML Results</caption>
<tr><th>Model</th><th>Score</th></tr>
<tr><td>A</td><td>0.91</td></tr>
</table>
""",
        doc_name="paper",
    )

    titles = {t.title for t in doc.tables}
    assert "HTML Results" in titles
    html_table = next(t for t in doc.tables if t.title == "HTML Results")
    assert html_table.headers == ["Model", "Score"]
    assert html_table.rows == [["A", "0.91"]]
    assert html_table.source_ref.startswith("paper:")
