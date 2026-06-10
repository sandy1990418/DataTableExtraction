from __future__ import annotations

import io

import pytest
from pptx import Presentation

from app.tools import table_pptx


def test_build_tables_pptx_creates_table_slide() -> None:
    data = table_pptx.build_tables_pptx(
        [{"title": "Results", "headers": ["Model", "Score"], "rows": [["A", "92.5"]]}]
    )

    presentation = Presentation(io.BytesIO(data))
    assert len(presentation.slides) == 1
    table = next(
        shape.table
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_table", False)
    )
    assert table.cell(1, 0).text == "A"
    assert table.cell(1, 1).text == "92.5"


def test_build_tables_pptx_rejects_empty_headers() -> None:
    with pytest.raises(ValueError, match="headers must not be empty"):
        table_pptx.build_tables_pptx([{"title": "Results", "headers": [], "rows": []}])


def test_safe_filename_preserves_unicode() -> None:
    assert table_pptx.safe_pptx_filename("中文 / 報表?.pptx") == "中文 報表?.pptx"
    assert table_pptx.ascii_fallback_filename("中文 / 報表?.pptx") == "table.pptx"
