from __future__ import annotations

import io

from fastapi.testclient import TestClient
from pptx import Presentation

from app.main import create_app
from app.models.data_table import (
    DataTableColumn,
    DataTableSchema,
    GroundedCell,
    GroundedDataTable,
    GroundedRow,
    RowEntity,
)
from app.routes import data_table as data_table_route
from app.tools import table_pptx


def test_app_exposes_only_data_table_business_routes() -> None:
    routes = {
        (method, route.path)
        for route in create_app().routes
        for method in getattr(route, "methods", [])
    }

    assert ("POST", "/data-table") in routes
    assert ("GET", "/download/{token}") in routes
    removed_paths = {"/analyze", "/chat", "/evidence", "/health", "/outline", "/render"}
    assert not any(path in removed_paths for _, path in routes)


def _make_table() -> GroundedDataTable:
    return GroundedDataTable(
        **{
            "schema": DataTableSchema(
                title="Exported Results",
                intent="Summarize results",
                columns=[
                    DataTableColumn(
                        name="Model",
                        role="entity",
                        description="Model name",
                        required=True,
                    ),
                    DataTableColumn(
                        name="Score",
                        role="metric",
                        description="Model score",
                        value_type="number",
                    ),
                ],
            ),
            "rows": [
                GroundedRow(
                    entity=RowEntity(entity_id="model-a", name="Model A"),
                    cells={
                        "Model": GroundedCell(value="Model A", status="supported"),
                        "Score": GroundedCell(value=92.5, status="supported"),
                    },
                )
            ],
            "metrics": {"row_count": 1, "column_count": 2},
        }
    )


def test_data_table_response_includes_downloadable_pptx(monkeypatch, tmp_path) -> None:
    async def fake_build_evidence(parsed_docs, settings, analyze_images):
        return []

    async def fake_generate_data_table(**kwargs):
        return _make_table()

    monkeypatch.setattr(data_table_route, "resolve_documents", lambda documents: [])
    monkeypatch.setattr(data_table_route, "build_evidence", fake_build_evidence)
    monkeypatch.setattr(data_table_route, "generate_data_table", fake_generate_data_table)
    monkeypatch.setattr(table_pptx, "STORE_DIR", tmp_path)

    client = TestClient(create_app())
    response = client.post(
        "/data-table",
        json={"documents": [{"content": "# Results"}], "hint": "Summarize results"},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["type"] == "data_table"
    assert payload["pptx"]["type"] == "download"
    assert payload["pptx"]["filename"] == "Exported Results.pptx"

    download = client.get(payload["pptx"]["url"])
    assert download.status_code == 200
    presentation = Presentation(io.BytesIO(download.content))
    assert len(presentation.slides) == 1
    table_shape = next(
        shape for shape in presentation.slides[0].shapes if getattr(shape, "has_table", False)
    )
    assert table_shape.table.cell(1, 0).text == "Model A"
    assert table_shape.table.cell(1, 1).text == "92.5"


def test_data_table_empty_evidence_flow_still_exports_pptx(monkeypatch, tmp_path) -> None:
    monkeypatch.setattr(table_pptx, "STORE_DIR", tmp_path)
    client = TestClient(create_app())

    response = client.post(
        "/data-table",
        json={"documents": [{"content": ""}], "analyze_images": False},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["table"] == {"headers": ["Entity"], "rows": []}
    assert payload["pptx"]["type"] == "download"
    assert client.get(payload["pptx"]["url"]).status_code == 200
