"""Build and temporarily store table-only PPTX files."""

from __future__ import annotations

import io
import re
import time
import uuid
from pathlib import Path
from urllib.parse import quote

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.4)
TITLE_H = Inches(0.7)
TITLE_TOP = Inches(0.2)
CONTENT_TOP = TITLE_TOP + TITLE_H + Inches(0.15)
CONTENT_H = SLIDE_H - CONTENT_TOP - MARGIN
ROW_H = Inches(0.38)
DATA_FONT = 11
MIN_FONT_PT = 7
MAX_ROWS_HARD = 60

HEADER_BG = RGBColor(0x2F, 0x54, 0x96)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ALT_BG = RGBColor(0xDD, 0xE8, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

STORE_DIR = Path("/tmp/backend-table-pptx")
_store: dict[str, tuple[Path, str, float]] = {}


def safe_pptx_filename(filename: str) -> str:
    clean = re.sub(r"[\\/\r\n\t]+", " ", filename).strip()
    clean = re.sub(r"\s+", " ", clean)
    clean = clean or "table.pptx"
    if not clean.lower().endswith(".pptx"):
        clean = f"{clean}.pptx"
    return clean


def ascii_fallback_filename(filename: str) -> str:
    clean = safe_pptx_filename(filename)
    ascii_name = re.sub(r"[^A-Za-z0-9._ -]+", "", clean).strip(" .")
    if ascii_name.lower() == "pptx":
        return "table.pptx"
    return ascii_name or "table.pptx"


def content_disposition(filename: str) -> str:
    safe_name = safe_pptx_filename(filename)
    fallback = ascii_fallback_filename(safe_name)
    return f'attachment; filename="{fallback}"; filename*=UTF-8\'\'{quote(safe_name)}'


def store_pptx(data: bytes, filename: str, ttl_seconds: int) -> str:
    token = uuid.uuid4().hex
    STORE_DIR.mkdir(parents=True, exist_ok=True)
    path = STORE_DIR / f"{token}.pptx"
    path.write_bytes(data)
    _store[token] = (path, safe_pptx_filename(filename), time.monotonic() + ttl_seconds)
    return token


def get_pptx_entry(token: str) -> tuple[bytes, str] | None:
    entry = _store.get(token)
    if entry is None:
        return None
    path, filename, expiry = entry
    if time.monotonic() > expiry or not path.exists():
        _store.pop(token, None)
        path.unlink(missing_ok=True)
        return None
    return path.read_bytes(), filename


def evict_expired() -> None:
    now = time.monotonic()
    expired = [token for token, entry in _store.items() if now > entry[2]]
    for token in expired:
        path, _, _ = _store.pop(token)
        path.unlink(missing_ok=True)


def _add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(MARGIN, TITLE_TOP, SLIDE_W - MARGIN * 2, TITLE_H)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text
    runs = frame.paragraphs[0].runs
    if runs:
        runs[0].font.size = Pt(24)
        runs[0].font.bold = True


def _rows_per_slide(font_pt: float) -> int:
    row_height = Inches(font_pt / 72 * 1.8)
    return max(1, int(CONTENT_H / row_height) - 1)


def _font_and_chunks(row_count: int) -> tuple[float, list[int]]:
    font_pt = DATA_FONT
    while font_pt >= MIN_FONT_PT:
        rows_per_slide = _rows_per_slide(font_pt)
        if rows_per_slide >= row_count:
            return font_pt, [row_count]
        font_pt -= 1

    rows_per_slide = min(_rows_per_slide(MIN_FONT_PT), MAX_ROWS_HARD)
    chunks = []
    remaining = row_count
    while remaining > 0:
        chunk_size = min(rows_per_slide, remaining)
        chunks.append(chunk_size)
        remaining -= chunk_size
    return MIN_FONT_PT, chunks or [0]


def _column_widths(headers: list[str], rows: list[list[str]], total_width: int) -> list[int]:
    max_lengths = [len(header) for header in headers]
    for row in rows:
        for index in range(len(headers)):
            value = row[index] if index < len(row) else ""
            max_lengths[index] = max(max_lengths[index], len(value))

    minimum = int(total_width / len(headers) * 0.4)
    total_chars = sum(max_lengths) or 1
    widths = [max(minimum, int(total_width * length / total_chars)) for length in max_lengths]
    if sum(widths) > total_width:
        scale = total_width / sum(widths)
        widths = [max(minimum, int(width * scale)) for width in widths]
    widths[-1] = max(minimum, total_width - sum(widths[:-1]))
    return widths


def _add_table(slide, headers: list[str], rows: list[list[str]], font_pt: float) -> None:
    width = SLIDE_W - MARGIN * 2
    height = min(CONTENT_H, ROW_H * (len(rows) + 1))
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        MARGIN,
        CONTENT_TOP,
        width,
        height,
    ).table

    for index, column_width in enumerate(_column_widths(headers, rows, int(width))):
        table.columns[index].width = column_width

    for index, text in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        runs = cell.text_frame.paragraphs[0].runs
        if runs:
            runs[0].font.bold = True
            runs[0].font.size = Pt(font_pt)
            runs[0].font.color.rgb = HEADER_FG

    for row_index, row in enumerate(rows):
        background = ALT_BG if row_index % 2 == 0 else WHITE
        for column_index in range(len(headers)):
            cell = table.cell(row_index + 1, column_index)
            cell.text = row[column_index] if column_index < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = background
            runs = cell.text_frame.paragraphs[0].runs
            if runs:
                runs[0].font.size = Pt(font_pt)


def _validate_table(headers: list[str], rows: list[list[str]]) -> None:
    if not headers:
        raise ValueError("headers must not be empty")
    if any(not header.strip() for header in headers):
        raise ValueError("headers must not contain blank values")
    if rows is None:
        raise ValueError("rows must not be None")


def build_tables_pptx(tables: list[dict]) -> bytes:
    if not tables:
        raise ValueError("tables must not be empty")

    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    blank_layout = presentation.slide_layouts[6]

    for raw_table in tables:
        title = str(raw_table.get("title") or "Data Table")
        headers = [str(header) for header in raw_table.get("headers", [])]
        rows = [[str(cell) for cell in row] for row in raw_table.get("rows", [])]
        _validate_table(headers, rows)

        font_pt, chunks = _font_and_chunks(len(rows))
        offset = 0
        for page_index, chunk_size in enumerate(chunks, 1):
            slide = presentation.slides.add_slide(blank_layout)
            page_title = title if len(chunks) == 1 else f"{title} ({page_index}/{len(chunks)})"
            _add_title(slide, page_title)
            chunk_rows = rows[offset : offset + chunk_size]
            offset += chunk_size
            _add_table(slide, headers, chunk_rows, font_pt)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
